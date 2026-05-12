"""
文档中间表示 —— 引擎与各格式适配器之间的统一契约。

P2-2 设计（冻结）：
  - TextPart: 文本片段 + 位置标识 + 人类可读展示标签
  - DocumentText: 格式无关的文档表示，引擎唯一天然输入
  - 各 from_xxx() 工厂方法由 io 层适配器实现

冻结范围（v0.3.0 承诺）：
  - TextPart 字段名: location, text, display
  - DocumentText 字段名: parts, full_text
  - 工厂方法: from_docx / from_text / from_markdown / from_html / from_jsonl / from_csv / from_pdf / from_pptx
  - gehd_check(text: DocumentText, ...) 签名（P2-2 实现）
"""

from __future__ import annotations

import re
from dataclasses import dataclass, field
from pathlib import Path


@dataclass(frozen=True)
class TextPart:
    """文档中的单个文本片段。

    Attributes:
        location: 机器可读的位置标识（如 "P1", "T2[3,1]"）
        text: 文本内容
        display: 人类可读的位置标签（如 "段落 1", "表格 2 行 3 列 1"）
    """

    location: str
    text: str

    @property
    def display(self) -> str:
        """生成人类可读的位置标签。

        转换规则：
          P1       → 段落 1
          T2[3,1]  → 表格 2 行 3 列 1
          其他      → 原样返回
        """
        m = re.match(r'^P(\d+)$', self.location)
        if m:
            return f'段落 {m.group(1)}'

        m = re.match(r'^T(\d+)\[(\d+),(\d+)\]$', self.location)
        if m:
            return f'表格 {m.group(1)} 行 {m.group(2)} 列 {m.group(3)}'

        return self.location


@dataclass(frozen=True)
class DocumentText:
    """格式无关的文档表示 —— 引擎的唯一天然输入。

    由各格式适配器负责构造：
      docx  → DocumentText.from_docx(path)
      txt   → DocumentText.from_text(path)
      md    → DocumentText.from_markdown(path)
      直接  → DocumentText(parts=(TextPart(...), ...))
    """

    parts: tuple[TextPart, ...]
    full_text: str = field(init=False)

    def __post_init__(self):
        object.__setattr__(self, 'full_text', '\n'.join(p.text for p in self.parts))

    # ---- 工厂方法 ----

    @classmethod
    def from_docx(cls, filepath: str | Path) -> DocumentText:
        """从 .docx 文件构造 DocumentText。

        内部调用现有 docx_reader + text_extractor，零破坏性变更。
        """
        filepath = Path(filepath)

        # 延迟导入，避免循环依赖
        from hallucination_checker.engine.extractors.text_extractor import (
            extract_all_text,
        )
        from hallucination_checker.io.docx_reader import load_docx

        doc = load_docx(str(filepath))
        raw_parts = extract_all_text(doc)
        parts = tuple(TextPart(location=loc, text=txt) for loc, txt in raw_parts)
        return cls(parts=parts)

    @classmethod
    def from_text(cls, filepath: str | Path) -> DocumentText:
        """从纯文本文件构造。每行一个 TextPart。"""
        filepath = Path(filepath)
        with open(filepath, encoding='utf-8') as f:
            lines = f.read().splitlines()

        parts = tuple(
            TextPart(location=f'L{i + 1}', text=line)
            for i, line in enumerate(lines)
        )
        return cls(parts=parts)

    @classmethod
    def from_markdown(cls, filepath: str | Path) -> DocumentText:
        """从 Markdown 文件构造。

        标题 (#...) 和普通段落各自一个 TextPart。
        空行分隔段落，连续非空行合并为一个段落。
        """
        import re

        filepath = Path(filepath)
        with open(filepath, encoding='utf-8') as f:
            content = f.read()

        parts: list[TextPart] = []
        loc_num = 0

        for block in content.split('\n\n'):
            block = block.strip()
            if not block:
                continue

            loc_num += 1
            if block.startswith('#'):
                m = re.match(r'^(#{1,6})\s+(.+)', block)
                if m:
                    level = len(m.group(1))
                    title_text = m.group(2)
                    parts.append(TextPart(
                        location=f'H{level}-{title_text[:20]}',
                        text=block,
                    ))
                else:
                    parts.append(TextPart(location=f'P-标题{loc_num}', text=block))
            else:
                parts.append(TextPart(location=f'P-段落{loc_num}', text=block))

        return cls(parts=tuple(parts))

    @classmethod
    def from_html(cls, filepath: str | Path) -> DocumentText:
        """从 HTML 文件构造。

        提取可见文本（忽略标签和 <script>/<style>），每个可见文本块一个 TextPart。
        """
        import re

        filepath = Path(filepath)
        with open(filepath, encoding='utf-8') as f:
            html = f.read()

        # 移除 script 和 style 内容
        html = re.sub(r'<script[^>]*>.*?</script>', '', html, flags=re.DOTALL | re.IGNORECASE)
        html = re.sub(r'<style[^>]*>.*?</style>', '', html, flags=re.DOTALL | re.IGNORECASE)

        # 提取可见文本片段
        text = re.sub(r'<[^>]+>', '\n', html)
        text = re.sub(r'&nbsp;', ' ', text)
        text = re.sub(r'&[a-z]+;', ' ', text)

        parts: list[TextPart] = []
        for i, line in enumerate(text.split('\n'), 1):
            line = line.strip()
            if line and len(line) >= 3:
                parts.append(TextPart(location=f'HTML-{i}', text=line))

        return cls(parts=tuple(parts))

    @classmethod
    def from_jsonl(cls, filepath: str | Path) -> DocumentText:
        """从 JSONL 文件构造。

        每行一个 JSON 对象，提取 text/content/body 字段，无文本字段则用整行原文。
        """
        import json

        filepath = Path(filepath)
        parts: list[TextPart] = []
        loc = 0

        with open(filepath, encoding='utf-8') as f:
            for raw in f:
                line = raw.strip()
                if not line:
                    continue
                loc += 1
                try:
                    obj = json.loads(line)
                    body = obj.get('text') or obj.get('content') or obj.get('body') or line
                except json.JSONDecodeError:
                    body = line
                parts.append(TextPart(location=f'JSONL-{loc}', text=str(body)))

        return cls(parts=tuple(parts))

    @classmethod
    def from_csv(cls, filepath: str | Path) -> DocumentText:
        """从 CSV 文件构造。

        每行一个 TextPart，列用 ' | ' 拼接。
        """
        import csv

        filepath = Path(filepath)
        parts: list[TextPart] = []
        loc = 0

        with open(filepath, encoding='utf-8-sig') as f:
            reader = csv.reader(f)
            for row in reader:
                loc += 1
                line = ' | '.join(cell.strip() for cell in row if cell.strip())
                if line:
                    parts.append(TextPart(location=f'CSV-{loc}', text=line))

        return cls(parts=tuple(parts))

    @classmethod
    def from_pdf(cls, filepath: str | Path) -> DocumentText:
        """从 PDF 文件构造。

        使用 pymupdf 逐页提取文本。pymupdf 不可用则抛 ImportError。
        """
        try:
            import fitz  # pymupdf
        except ImportError:
            raise ImportError('pip install pymupdf 以支持 PDF 格式') from None

        filepath = Path(filepath)
        try:
            doc = fitz.open(str(filepath))
        except (ValueError, OSError, RuntimeError) as e:
            raise ValueError(f'无法打开 PDF 文件（可能已损坏）: {filepath.name}') from e
        parts: list[TextPart] = []

        try:
            for i, page in enumerate(doc, 1):
                text = page.get_text().strip()
                if text:
                    parts.append(TextPart(location=f'PDF-页{i}', text=text))
        finally:
            doc.close()

        return cls(parts=tuple(parts))

    @classmethod
    def from_pptx(cls, filepath: str | Path) -> DocumentText:
        """从 PPTX 文件构造。

        使用 python-pptx 逐幻灯片提取文本。python-pptx 不可用则抛 ImportError。
        """
        try:
            from pptx import Presentation
            from pptx.exc import PackageNotFoundError
        except ImportError:
            raise ImportError('pip install python-pptx 以支持 PPTX 格式') from None

        filepath = Path(filepath)
        try:
            prs = Presentation(str(filepath))
        except (PackageNotFoundError, ValueError, OSError) as e:
            raise ValueError(f'无法打开 PPTX 文件（可能已损坏）: {filepath.name}') from e
        parts: list[TextPart] = []

        try:
            for i, slide in enumerate(prs.slides, 1):
                texts: list[str] = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                        if t:
                            texts.append(t)
                if texts:
                    parts.append(TextPart(location=f'PPTX-幻灯片{i}', text='\n'.join(texts)))
        except RuntimeError as e:
            raise ValueError(f'PPTX 文件处理失败: {filepath.name}') from e

        return cls(parts=tuple(parts))
