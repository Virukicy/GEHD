"""
多格式适配器单元测试 — HTML / JSONL / CSV / PDF / PPTX。
"""

import csv
import json
import os
import sys
import tempfile
from pathlib import Path

sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', 'src'))

import pytest

from hallucination_checker.io.document_text import DocumentText


class TestHTML:
    def test_basic_html(self):
        with tempfile.NamedTemporaryFile(
            mode='w', suffix='.html', delete=False, encoding='utf-8',
        ) as f:
            f.write('<html><body><p>段落一</p><p>段落二</p><script>alert(1)</script></body></html>')
            tmp = f.name
        try:
            doc = DocumentText.from_html(tmp)
            assert len(doc.parts) >= 2
            assert any('段落一' in p.text for p in doc.parts)
        finally:
            Path(tmp).unlink(missing_ok=True)


class TestJSONL:
    def test_basic_jsonl(self):
        with tempfile.NamedTemporaryFile(
            mode='w', suffix='.jsonl', delete=False, encoding='utf-8',
        ) as f:
            f.write(json.dumps({'text': '记录一'}, ensure_ascii=False) + '\n')
            f.write(json.dumps({'content': '记录二'}, ensure_ascii=False) + '\n')
            tmp = f.name
        try:
            doc = DocumentText.from_jsonl(tmp)
            assert len(doc.parts) == 2
            assert '记录一' in doc.parts[0].text
        finally:
            Path(tmp).unlink(missing_ok=True)


class TestCSV:
    def test_basic_csv(self):
        with tempfile.NamedTemporaryFile(
            mode='w', suffix='.csv', delete=False, encoding='utf-8-sig',
        ) as f:
            writer = csv.writer(f)
            writer.writerow(['列A', '列B'])
            writer.writerow(['值1', '值2'])
            tmp = f.name
        try:
            doc = DocumentText.from_csv(tmp)
            assert len(doc.parts) >= 1
            assert '值1' in doc.full_text
        finally:
            Path(tmp).unlink(missing_ok=True)


class TestPDF:
    def test_basic_pdf(self):
        try:
            import fitz  # noqa: F401
        except ImportError:
            pytest.skip('pymupdf 未安装')
        # 创建最小的合法 PDF
        tmp = Path(tempfile.mkdtemp()) / 'test.pdf'
        try:
            doc_pdf = fitz.open()
            page = doc_pdf.new_page()
            page.insert_text((72, 72), 'Hello PDF Test')
            doc_pdf.save(str(tmp))
            doc_pdf.close()

            doc = DocumentText.from_pdf(tmp)
            assert len(doc.parts) >= 1
            assert 'Hello' in doc.full_text
        finally:
            tmp.unlink(missing_ok=True)

    def test_corrupted_pdf(self):
        with tempfile.NamedTemporaryFile(
            mode='w', suffix='.pdf', delete=False, encoding='utf-8',
        ) as f:
            f.write('这是一个损坏的PDF文件')
            tmp = f.name
        try:
            with pytest.raises(ValueError, match='无法打开|已损坏'):
                DocumentText.from_pdf(tmp)
        finally:
            Path(tmp).unlink(missing_ok=True)


class TestPPTX:
    def test_basic_pptx(self):
        try:
            from pptx import Presentation  # noqa: F401
        except ImportError:
            pytest.skip('python-pptx 未安装')
        tmp = Path(tempfile.mkdtemp()) / 'test.pptx'
        try:
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = '演示文稿标题'
            prs.save(str(tmp))

            doc = DocumentText.from_pptx(tmp)
            assert len(doc.parts) >= 1
            assert '演示文稿标题' in doc.full_text
        finally:
            tmp.unlink(missing_ok=True)

    def test_corrupted_pptx(self):
        with tempfile.NamedTemporaryFile(
            mode='w', suffix='.pptx', delete=False, encoding='utf-8',
        ) as f:
            f.write('这是一个损坏的PPTX文件')
            tmp = f.name
        try:
            with pytest.raises(ValueError, match='无法打开|已损坏'):
                DocumentText.from_pptx(tmp)
        finally:
            Path(tmp).unlink(missing_ok=True)
